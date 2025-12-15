# ✅ Premium Study & Coding Document Engine (FINAL – RULE RESPECT)
# Dev — @lakshitpatidar
# NOTHING REMOVED • UI SAME • FEATURES SAME

import os
import re
import tempfile

from pyrogram import Client, filters
from pyrogram.types import Message

import pdfplumber
from docx import Document as DocxDocument
from pptx import Presentation

# ================= CONFIG =================
API_ID = 12767104
API_HASH = "a0ce1daccf78234927eb68a62f894b97"
BOT_TOKEN = "8412966210:AAFahj1GR4Jrrf_MVrgvkqks3iFK2IVsR0k"

MAX_LEN = 3500  # mono safe

# ================= CODE FILE TYPES =================
CODE_EXTENSIONS = {
    ".py", ".js", ".ts", ".html", ".css", ".json",
    ".java", ".cpp", ".c", ".cs", ".php", ".go",
    ".rs", ".sh", ".yaml", ".yml"
}

# ================= START TEXT =================
START_TEXT = (
    "Hey There 👋\n"
    "I Am 𝕷𝖆𝖐𝖘𝖍𝖎𝖙 𝕻𝖆𝖙𝖎𝖉𝖆𝖗\n\n"
    "✦ Study & Coding Document Engine\n"
    "✦ Premium Extract • Image First\n"
    "✦ Legend Fonts • Clean UI\n\n"
    "━━━━━━━━━━━━━━━━━━━━━━\n"
    "Dev — @lakshitpatidar"
)

# ================= APP =================
app = Client(
    "premium_doc_bot",
    api_id=API_ID,
    api_hash=API_HASH,
    bot_token=BOT_TOKEN
)

# ================= SAFE CLEAN (ONLY REAL JUNK) =================
def clean_only_junk(text: str) -> str:
    """
    IMPORTANT:
    - Does NOT remove Q1, Ans, headings, labels if present in file
    - Preserves spacing & blank lines
    - Removes ONLY absolute junk not part of content
    """
    out = []

    for raw in text.splitlines():
        line = raw.rstrip("\n")

        # keep blank lines (spacing)
        if not line.strip():
            out.append("")
            continue

        stripped = line.strip()

        # pure symbols line (design junk)
        if re.fullmatch(r"[•➤▶▪◆■★☆━─_]+", stripped):
            continue

        # random alphabet garbage (abcdef, qwertyui)
        if re.fullmatch(r"[a-zA-Z]{6,}", stripped):
            continue

        # consonant OCR garbage
        if re.fullmatch(r"[bcdfghjklmnpqrstvwxyz]{7,}", stripped.lower()):
            continue

        # vertical spaced OCR junk (H U C S T)
        if re.fullmatch(r"(?:[A-Za-z]\s+){3,}[A-Za-z]?", stripped):
            continue

        out.append(line)

    return "\n".join(out)

# ================= MONO SEND (HARD SAFE) =================
async def send_all_mono(msg: Message, text: str, lang: str = ""):
    chunks = [text[i:i + MAX_LEN] for i in range(0, len(text), MAX_LEN)]
    for chunk in chunks:
        await msg.reply_text(f"```{lang}\n{chunk.rstrip()}\n```")

# ================= PDF EXTRACT =================
def extract_pdf(path, img_dir):
    text, images = [], []
    with pdfplumber.open(path) as pdf:
        pages = len(pdf.pages)
        for page in pdf.pages:
            t = page.extract_text(x_tolerance=2, y_tolerance=2)
            if t:
                text.append(t)
            for i, img in enumerate(page.images):
                try:
                    im = page.crop(
                        (img["x0"], img["top"], img["x1"], img["bottom"])
                    ).to_image(resolution=150)
                    p = os.path.join(img_dir, f"pdf_{len(images)}.png")
                    im.save(p)
                    images.append(p)
                except:
                    pass
    return "\n\n".join(text), pages, len(images), images

# ================= DOCX EXTRACT =================
def extract_docx(path, img_dir):
    doc = DocxDocument(path)
    text, images = [], []

    for p in doc.paragraphs:
        text.append(p.text)

    for rel in doc.part._rels.values():
        if "image" in rel.reltype:
            p = os.path.join(img_dir, f"docx_{len(images)}.png")
            with open(p, "wb") as f:
                f.write(rel.target_part.blob)
            images.append(p)

    return "\n\n".join(text), len(doc.paragraphs), len(images), images

# ================= PPTX EXTRACT =================
def extract_pptx(path, img_dir):
    prs = Presentation(path)
    text = []

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)

    return "\n\n".join(text), len(prs.slides), 0, []

# ================= START =================
@app.on_message(filters.command("start"))
async def start_handler(_, msg: Message):
    await msg.reply_text(START_TEXT)

# ================= DOCUMENT HANDLER =================
@app.on_message(filters.document)
async def document_handler(_, msg: Message):
    await msg.reply_text("✦ Processing document…")

    with tempfile.TemporaryDirectory() as tmp:
        file_path = await msg.download(
            file_name=os.path.join(tmp, msg.document.file_name)
        )

        ext = os.path.splitext(file_path)[1].lower()
        is_code = ext in CODE_EXTENSIONS

        # -------- CODE FILES (RAW ONLY) --------
        if is_code:
            with open(file_path, "r", errors="ignore") as f:
                text = f.read()

            await msg.reply_text(
                "🧠 𝗖𝗢𝗗𝗘 𝗘𝗫𝗧𝗥𝗔𝗖𝗧\n"
                "➤ Raw Source • Exact Spacing • Copy Ready"
            )
            await send_all_mono(msg, text, ext.replace(".", ""))
            await msg.reply_text("━━━━━━━━━━━━━━━━━━━━━━\nDev — @lakshitpatidar")
            return

        # -------- DOCUMENT FILES --------
        img_dir = os.path.join(tmp, "images")
        os.makedirs(img_dir, exist_ok=True)

        if ext == ".pdf":
            text, pages, img_count, images = extract_pdf(file_path, img_dir)
        elif ext == ".docx":
            text, pages, img_count, images = extract_docx(file_path, img_dir)
        elif ext == ".pptx":
            text, pages, img_count, images = extract_pptx(file_path, img_dir)
        else:
            with open(file_path, "r", errors="ignore") as f:
                text = f.read()
            pages, img_count, images = 1, 0, []

        await msg.reply_text(
            "📄 DOCUMENT INFO\n"
            f"✦ Source : {msg.document.file_name}\n"
            f"✦ Pages : {pages}\n"
            f"✦ Images : {img_count}\n"
            "━━━━━━━━━━━━━━━━━━━━━━"
        )

        for img in images:
            await msg.reply_photo(img)

        # CLEAN ONLY REAL JUNK (NOT CONTENT)
        text = clean_only_junk(text)
        if not text.strip():
            text = "⟡ No meaningful text found."

        await msg.reply_text(
            "📘 𝗦𝗧𝗨𝗗𝗬 𝗘𝗫𝗧𝗥𝗔𝗖𝗧\n"
            "➤ File-Accurate • Original Spacing • No Fabrication"
        )
        await send_all_mono(msg, text)

        await msg.reply_text("━━━━━━━━━━━━━━━━━━━━━━\nDev — @lakshitpatidar")

# ================= RUN =================
app.run()
